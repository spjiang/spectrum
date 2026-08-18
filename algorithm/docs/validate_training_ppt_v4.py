#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""校验高光谱算法体系 v4 PPT 的结构、口径与基础版面。"""

from __future__ import annotations

from collections import Counter
from pathlib import Path

from pptx import Presentation

from generate_training_ppt_v4 import IMPL, LIST_MD, OUT, parse_algorithms


ALLOWED_FONTS = {"Noto Sans SC", "Inter", "PingFang SC"}
STALE_PHRASES = {"已实现 9 项", "其余 36 项", "已实现 11 项", "其余 34 项"}


def slide_text(slide) -> str:
    return "\n".join(
        shape.text.strip()
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    )


def validate() -> None:
    assert OUT.exists(), f"未找到输出文件：{OUT}"
    prs = Presentation(OUT)
    algorithms = parse_algorithms(LIST_MD.read_text(encoding="utf-8"))

    assert len(prs.slides) == 62, f"页数异常：{len(prs.slides)}"
    all_text = "\n".join(slide_text(slide) for slide in prs.slides)

    for item in algorithms:
        expected = f"{item['num']:02d}. {item['name']}"
        assert expected in all_text, f"缺少算法标题：{expected}"

    implemented_slides = 0
    skeleton_slides = 0
    font_counter: Counter[str] = Counter()
    undersized = []
    out_of_bounds = []
    invalid_titles = []

    for slide_index, slide in enumerate(prs.slides, 1):
        text = slide_text(slide)
        if "可运行｜教学级实现" in text:
            implemented_slides += 1
        if "骨架｜核心算法待实现" in text:
            skeleton_slides += 1

        nonempty_texts = []
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                out_of_bounds.append((slide_index, shape.name))
            if not getattr(shape, "has_text_frame", False):
                continue
            shape_text = shape.text.strip()
            if shape_text:
                nonempty_texts.append(shape_text)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    if run.font.name:
                        font_counter[run.font.name] += 1
                    if run.font.size and run.font.size.pt < 11:
                        undersized.append((slide_index, run.text[:30], run.font.size.pt))

        first = nonempty_texts[0] if nonempty_texts else ""
        if first and (first.isdigit() or first.startswith("#")):
            invalid_titles.append((slide_index, first))

    assert implemented_slides == len(IMPL) == 12, f"可运行页数异常：{implemented_slides}"
    assert skeleton_slides == 33, f"骨架页数异常：{skeleton_slides}"
    assert not out_of_bounds, f"存在越界形状：{out_of_bounds[:10]}"
    assert not invalid_titles, f"存在残缺标题：{invalid_titles}"
    assert not undersized, f"存在小于 11pt 的正文：{undersized[:10]}"
    assert not (set(font_counter) - ALLOWED_FONTS), f"存在非标准字体：{set(font_counter) - ALLOWED_FONTS}"
    for phrase in STALE_PHRASES:
        assert phrase not in all_text, f"存在过期口径：{phrase}"

    print(
        "PPT 校验通过：62 页；45 项算法；"
        f"{implemented_slides} 可运行；{skeleton_slides} 骨架；"
        f"0 个越界形状；字体 {dict(font_counter)}"
    )


if __name__ == "__main__":
    validate()
