#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""使用 Pillow 生成关键幻灯片预览，供无 Office 渲染环境下抽检。"""

from __future__ import annotations

import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN


ROOT = Path(__file__).resolve().parent
PPTX_PATH = ROOT / "高光谱算法体系-v4.pptx"
OUTPUT_DIR = ROOT / "ppt-v4-preview"
SLIDES = {
    1: "01-cover.png",
    4: "04-l0-l4-architecture.png",
    9: "09-algorithm-overview.png",
    11: "11-algorithm-skeleton.png",
    22: "22-algorithm-runnable.png",
    58: "58-service-capabilities.png",
    61: "61-summary.png",
}

WIDTH, HEIGHT = 1600, 900
FONT_CANDIDATES = [
    "/System/Library/Fonts/PingFang.ttc",
    "/System/Library/Fonts/STHeiti Medium.ttc",
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
]


def rgb_value(color, fallback=(255, 255, 255)):
    try:
        value = color.rgb
        if value is not None:
            return tuple(value)
    except (AttributeError, TypeError):
        pass
    return fallback


def font_path() -> str:
    for candidate in FONT_CANDIDATES:
        if Path(candidate).exists():
            return candidate
    raise FileNotFoundError("未找到可用中文字体")


def load_font(size: int):
    return ImageFont.truetype(font_path(), max(8, size))


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font, max_width: int) -> list[str]:
    if not text:
        return [""]
    lines = []
    for raw_line in text.splitlines() or [""]:
        if not raw_line:
            lines.append("")
            continue
        current = ""
        for char in raw_line:
            candidate = current + char
            if current and draw.textlength(candidate, font=font) > max_width:
                lines.append(current)
                current = char
            else:
                current = candidate
        lines.append(current)
    return lines


def shape_box(shape, prs):
    sx = WIDTH / prs.slide_width
    sy = HEIGHT / prs.slide_height
    return (
        int(shape.left * sx),
        int(shape.top * sy),
        int((shape.left + shape.width) * sx),
        int((shape.top + shape.height) * sy),
    )


def draw_shape_background(draw: ImageDraw.ImageDraw, shape, box) -> None:
    fill = None
    try:
        if shape.fill.type == MSO_FILL.SOLID:
            fill = rgb_value(shape.fill.fore_color)
    except (AttributeError, TypeError):
        pass
    if fill is None:
        return
    outline = None
    try:
        outline = rgb_value(shape.line.color, None)
    except (AttributeError, TypeError):
        pass
    rounded = getattr(shape, "auto_shape_type", None) == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    if rounded:
        draw.rounded_rectangle(box, radius=10, fill=fill, outline=outline, width=1)
    else:
        draw.rectangle(box, fill=fill, outline=outline, width=1)


def draw_shape_text(draw: ImageDraw.ImageDraw, shape, box) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    x1, y1, x2, y2 = box
    padding_x, padding_y = 10, 7
    cursor_y = y1 + padding_y
    max_width = max(10, x2 - x1 - padding_x * 2)
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text
        if not text:
            cursor_y += 4
            continue
        first_run = next((run for run in paragraph.runs if run.text), None)
        pt = first_run.font.size.pt if first_run and first_run.font.size else 14
        pixel_size = int(pt * WIDTH / 13.333 / 72)
        font = load_font(pixel_size)
        color = (20, 35, 55)
        if first_run:
            try:
                color = rgb_value(first_run.font.color, color)
            except (AttributeError, TypeError):
                pass
        lines = wrap_text(draw, text, font, max_width)
        line_height = int(pixel_size * 1.28)
        for line in lines:
            if cursor_y + line_height > y2:
                break
            line_width = draw.textlength(line, font=font)
            if paragraph.alignment == PP_ALIGN.CENTER:
                x = x1 + (x2 - x1 - line_width) / 2
            elif paragraph.alignment == PP_ALIGN.RIGHT:
                x = x2 - padding_x - line_width
            else:
                x = x1 + padding_x
            draw.text((x, cursor_y), line, font=font, fill=color)
            cursor_y += line_height
        cursor_y += max(2, int(pixel_size * 0.15))


def render_slide(prs: Presentation, slide_number: int, output: Path) -> None:
    image = Image.new("RGB", (WIDTH, HEIGHT), (247, 249, 252))
    draw = ImageDraw.Draw(image)
    slide = prs.slides[slide_number - 1]
    for shape in slide.shapes:
        box = shape_box(shape, prs)
        draw_shape_background(draw, shape, box)
        draw_shape_text(draw, shape, box)
    image.save(output, quality=95)


def main() -> None:
    prs = Presentation(PPTX_PATH)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    for slide_number, filename in SLIDES.items():
        render_slide(prs, slide_number, OUTPUT_DIR / filename)
    print(f"预览生成完成：{len(SLIDES)} 张，目录 {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
