#!/usr/bin/env python3
"""生成政企科技蓝风格的高光谱算法体系 v4 汇报 PPT。"""
from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
LIST_MD = ROOT / "docs" / "采集到算法-算法清单.md"
OUT = ROOT / "docs" / "高光谱算法体系-v4.pptx"

C = {
    "bg": (247, 249, 252),
    "ink": (20, 35, 55),
    "muted": (82, 99, 117),
    "line": (205, 216, 228),
    "card": (255, 255, 255),
    "primary": (24, 68, 125),
    "primary_dark": (14, 44, 85),
    "accent": (0, 154, 166),
    "warm": (209, 113, 29),
    "ok": (0, 135, 110),
    "white": (255, 255, 255),
    "soft_bg": (237, 246, 250),
}

FONT = "Noto Sans SC"
IMPL = {12, 20, 21, 22, 23, 27, 28, 34, 37, 40, 42, 45}
ALG_IDS = [
    "01_flight_planning", "02_sync_timestamp", "03_pos_solution", "04_flight_qc", "05_cloud_shadow",
    "06_dark_current", "07_bad_pixel", "08_destriping", "09_smile_keystone", "10_radiance_calibration",
    "11_relative_radiometric", "12_panel_reflectance", "13_atmospheric_correction", "14_brdf_correction",
    "15_geo_locate", "16_orthorectify", "17_mosaic", "18_color_balance", "19_multi_source_register",
    "20_bad_band_remove", "21_savgol_smooth", "22_normalize", "23_pca", "24_band_select",
    "25_superpixel", "26_patch_build", "27_ndvi", "28_ndre", "29_evi_savi",
    "30_ndmi_ndwi", "31_red_edge_params", "32_regression_inversion", "33_physical_inversion",
    "34_svm_rf_classify", "35_spectral_matching", "36_cnn1d_classify", "37_cnn3d_classify",
    "38_transformer_classify", "39_few_shot_classify", "40_detect_segment", "41_unmixing",
    "42_anomaly_detect", "43_change_detect", "44_postprocess_smooth", "45_parcel_zonal_stats",
]

# 正式补充：方法要点 / 业务价值 / 缺失影响 / 流程位置 / 关联 / 要点提示
ENRICH = {
    1: {
        "method": "依据测区边界、DEM 与相机视场，计算航高、旁向/航向重叠率与航迹间距；输出可执行航点序列。",
        "value": "保证测区完整覆盖与可镶嵌性，是后续几何处理与面积统计的前提条件。",
        "risk": "重叠不足将产生空洞；设计过密则增加架次与存储成本，延误交付周期。",
        "pipeline": "上机前规划 → 飞行采集 → #4 架次质检",
        "related": "关联：#3 POS、#16 正射、#17 镶嵌",
        "tip": "典型旁向重叠 60%–80%；分辨率由航高与像元尺寸共同决定。",
    },
    2: {
        "method": "以统一时钟或硬件触发脉冲为基准，对齐高光谱、RGB、POS 各传感器时间戳，建立帧级对应关系。",
        "value": "保障多源数据时空一致，避免几何校正阶段图—姿态错位。",
        "risk": "时间未对齐将导致正射漂移、镶嵌错位，后续排查成本高。",
        "pipeline": "多传感器采集 → 时间对齐 → #3 POS 解算 / 几何处理",
        "related": "关联：#3 POS、#15 地理定位、#19 多源配准",
        "tip": "工程上常输出「帧号—时间—姿态」对应表，作为几何链路输入。",
    },
    3: {
        "method": "融合 GPS 轨迹与 IMU 角速度/加速度（可选差分基站），解算逐帧位置 (x,y,z) 与姿态角。",
        "value": "为地理定位与正射校正提供外方位元素，使影像具备地图坐标基础。",
        "risk": "无可靠 POS 则无法正射与量测面积，数据仅能作为非地理影像使用。",
        "pipeline": "L0 原始包 → POS 解算 → #15/#16 几何校正",
        "related": "关联：#2 时间对齐、#15 地理定位、#16 正射",
        "tip": "差分 GPS（RTK/PPK）可显著降低平面误差，利于地块级量测。",
    },
    4: {
        "method": "检查丢帧、过曝/欠曝比例、POS 完整性与曝光日志异常，生成质检报告与重飞建议。",
        "value": "在处理管线入口拦截不可用架次，控制无效算力投入与返工风险。",
        "risk": "劣质数据进入后续环节，将造成指数/分类结果不可信并浪费处理资源。",
        "pipeline": "落地 → 质检闸门 → 通过则进入 L0→L1；不通过则补飞",
        "related": "关联：#1 航线规划、#5 云影检测、#10 辐射定标",
        "tip": "建议将饱和像元比例、丢帧率设为硬阈值，写入项目验收标准。",
    },
    5: {
        "method": "基于光谱/亮度阈值或分类模型识别云与阴影，输出二值或概率掩膜。",
        "value": "避免将遮挡区误判为作物或异常长势，提高专题图可信度。",
        "risk": "未掩膜时，指数出现虚假低值，分类产生伪斑块，影响农事判断。",
        "pipeline": "L0/L1 预览 → 云影掩膜 → 在 L3 分析中屏蔽或触发补采",
        "related": "关联：#4 架次质检、#27–30 指数、#34–37 分类",
        "tip": "卫星与高空数据更关键；低空亦需关注树影与薄云。",
    },
    6: {
        "method": "采集或调用暗电流参考帧，按像元/波段从原始 DN 中扣除本底：DN' = DN − Dark。",
        "value": "降低固定偏差，使弱信号区更接近真实辐射响应。",
        "risk": "未校正时本底抬高整体亮度，定量指数与定标系数产生系统偏差。",
        "pipeline": "原始 DN → 暗电流校正 → #7/#8 缺陷修复 → #10 辐射定标",
        "related": "关联：#7 坏像元、#8 去条带、#10 辐射定标",
        "tip": "暗帧应与工作积分时间、温度条件匹配，否则校正不完全。",
    },
    7: {
        "method": "依据坏像元表识别坏点/坏列，采用邻域插值或中值填充进行修复。",
        "value": "消除固定条纹与黑点伪影，避免被后续模型当作地物特征。",
        "risk": "坏列残留会在分类图中形成条状误差，降低产品合格率。",
        "pipeline": "暗电流后 → 坏像元修复 → #8 去条带或 #10 定标",
        "related": "关联：#6 暗电流、#8 去条带",
        "tip": "推扫探测器老化时坏列增多，建议纳入设备维保监测。",
    },
    8: {
        "method": "针对推扫方向周期性亮暗条纹，采用矩匹配、滤波或频域方法抑制条带噪声。",
        "value": "提升辐射一致性，改善镶嵌观感与分类稳定性。",
        "risk": "残留条带导致镶嵌接缝明显，分类沿飞行方向出现系统性误差。",
        "pipeline": "DN/辐亮度 → 去条带 → #10 辐射定标或 #12 反射率定标",
        "related": "关联：#7 坏像元、#10/#12 定标、#17 镶嵌",
        "tip": "机载/星载推扫高光谱的常见刚需步骤。",
    },
    9: {
        "method": "基于实验室光谱定标与传感器模型，校正 Smile（光谱维弯曲）与 Keystone（空间—光谱耦合）。",
        "value": "保证「同一波段对应同一波长」，支撑光谱库匹配与精密定量。",
        "risk": "未校正时与标准光谱库无法对齐，SAM/精细反演精度显著下降。",
        "pipeline": "传感器级预处理 → 光谱几何校正 → #35 光谱匹配 / #32–33 反演",
        "related": "关联：#10 辐射定标、#35 SAM、#31 红边参数",
        "tip": "农情粗指数有时可简化；矿物识别与精密反演建议保留。",
    },
    10: {
        "method": "应用定标系数与积分时间等元数据：L = gain × DN + offset，得到有物理单位的辐亮度。",
        "value": "实现跨传感器、跨架次的物理可比，是定量遥感的基础里程碑。",
        "risk": "停留在 DN 无法开展严谨定量对比，后续大气/反射率链路缺少物理输入。",
        "pipeline": "校正后 DN → L1 辐亮度 → #12 白板定标或 #13 大气校正",
        "related": "关联：#6–9 预处理、#12/#13 反射率",
        "tip": "输出为 L1 Radiance Cube/条带，需保留单位与波段中心波长元数据。",
    },
    11: {
        "method": "利用重叠区或伪不变特征，对多景进行相对辐射归一，消除增益与光照差异。",
        "value": "支持多架次/多时相公平对比，以及镶嵌前匀光。",
        "risk": "未归一时时间序列出现虚假「长势跳变」，告警规则失效。",
        "pipeline": "多景 L1/L2 → 相对辐射归一 → #17 镶嵌 / #43 变化检测",
        "related": "关联：#17 镶嵌、#18 匀色、#43 变化检测",
        "tip": "监测类项目应将辐射一致性纳入质量指标。",
    },
    12: {
        "method": "经验线法：利用地面白板/灰板同步测量，将辐亮度线性转换为近似地表反射率。",
        "value": "无人机农情最常用的反射率获取路径，为指数与分类提供标准底图。",
        "risk": "量纲错误将导致跨天不可比，指数与分类结果缺乏一致性。",
        "pipeline": "L1 辐亮度 + 白板 → 反射率立方体 → L2 分析入口",
        "related": "关联：#13 大气校正（替代路径）、#27–34 分析",
        "tip": "本服务已提供教学级可运行接口；白板应置于测区代表性光照条件下。",
    },
    13: {
        "method": "基于辐射传输（如 6S/MODTRAN 类）或暗目标法等，扣除大气吸收与散射。",
        "value": "获得更接近真实的地表反射率，支撑跨区域、跨季节定量产品。",
        "risk": "未校正时反射率含大气「染色」，区域对比与反演偏差增大。",
        "pipeline": "L1 辐亮度 + 大气参数 → 地表反射率（L2A 类）",
        "related": "关联：#12 白板定标（低空常用替代）、#32–33 反演",
        "tip": "卫星/高空项目更依赖；低空无人机多以白板法简化。",
    },
    14: {
        "method": "根据太阳—观测几何，对双向反射分布（BRDF）效应进行归一或校正。",
        "value": "减弱航带边缘「一边亮一边暗」，提高大范围合成一致性。",
        "risk": "角度效应残留会造成假异常斑与镶嵌花斑。",
        "pipeline": "反射率 + 角度元数据 → 角度归一反射率 → 镶嵌/监测",
        "related": "关联：#17 镶嵌、#11 相对辐射、#43 变化检测",
        "tip": "宽视场航测、多日多角度合景时优先考虑。",
    },
    15: {
        "method": "结合 POS 与相机内参，将条带像素投影到初始地理坐标系。",
        "value": "快速确认飞行落点，为精正射提供粗定位起点。",
        "risk": "无法预览落点时，任务质检与正射排错效率显著下降。",
        "pipeline": "影像条带 + POS → 粗地理参考 → #16 正射",
        "related": "关联：#3 POS、#16 正射、#4 质检",
        "tip": "精度通常不足以直接量亩，但满足任务级预览。",
    },
    16: {
        "method": "引入 DEM 与相机模型，消除地形起伏与姿态引起的几何畸变，生成正射条带。",
        "value": "使像素可准确叠合地块边界，支撑面积量测与 GIS 分析。",
        "risk": "几何不准导致地块错位，面积统计系统性失真。",
        "pipeline": "粗定位条带 + DEM → 正射条带 → #17 镶嵌",
        "related": "关联：#15 地理定位、#17 镶嵌、#45 地块统计",
        "tip": "「能否进入 GIS 与量亩」的关键分水岭。",
    },
    17: {
        "method": "在重叠区进行影像匹配与接缝线选择，将多航带合成为整景立方体。",
        "value": "形成测区完整 L2 产品，供整村/整园尺度分析与制图。",
        "risk": "仅保留碎片条带时，无法开展区域级专题产品。",
        "pipeline": "多条正射条带 → 镶嵌 → 整景反射率正射立方体（典型 L2）",
        "related": "关联：#16 正射、#18 匀色、#19 配准",
        "tip": "算法研发侧最常见的标准输入形态。",
    },
    18: {
        "method": "对镶嵌接缝进行匀光匀色与接缝线优化，消除色差硬边。",
        "value": "提升交付图可读性与验收通过率，不改变几何框架。",
        "risk": "接缝明显会降低产品观感，影响业务验收评价。",
        "pipeline": "初镶嵌 → 匀色/接缝优化 → 交付或进入分析",
        "related": "关联：#17 镶嵌、#11 相对辐射",
        "tip": "偏交付质量环节，但对项目验收影响直接。",
    },
    19: {
        "method": "将高光谱、RGB 正射与地块矢量配准到同一网格/坐标系。",
        "value": "支撑「真彩色目视 + 光谱算法 + 按地块统计」的标准业务形态。",
        "risk": "未配准导致标注错位、裁剪错误，面积与专题统计失效。",
        "pipeline": "HSI + RGB + shp/geojson → 共配准多层 → 标注/统计",
        "related": "关联：#17 镶嵌、#34–40 分析、#45 汇总",
        "tip": "打通「图、谱、田」三层数据是产品闭环关键。",
    },
    20: {
        "method": "按信噪比/水汽吸收特征剔除坏波段，并可做光谱维轻度去噪。",
        "value": "减少噪声波段对模型与指数的干扰，提高结果稳定性。",
        "risk": "保留强噪声波段会导致分类抖动与指数异常尖峰。",
        "pipeline": "L2 Cube → 坏波段剔除 → #21–23 特征处理或直接分析",
        "related": "关联：#21 SG、#23 PCA、#27–37 分析",
        "tip": "约 1400/1900 nm 水汽带常剔除；本服务可运行。",
    },
    21: {
        "method": "Savitzky–Golay 多项式滑动平滑；可选 Continuum Removal 突出吸收特征。",
        "value": "抑制光谱毛刺并增强吸收谷，利于精细识别与特征工程。",
        "risk": "噪声过大时弱特征被掩盖，匹配与反演灵敏度下降。",
        "pipeline": "清洗后光谱/Cube → 平滑/去包络 → 匹配、回归或深度学习",
        "related": "关联：#20 去噪、#35 SAM、#32 回归",
        "tip": "窗口长度与多项式阶次需匹配波段采样密度；本服务可运行。",
    },
    22: {
        "method": "对光谱向量或立方体做 z-score / min-max 等变换，统一量纲。",
        "value": "改善机器学习收敛与跨相机迁移能力。",
        "risk": "未标准化时模型偏向数值幅度大的波段，泛化变差。",
        "pipeline": "Cube → 标准化特征 → #23 PCA / #34–38 分类",
        "related": "关联：#23 降维、#34 SVM、#36–38 深度学习",
        "tip": "几乎所有 ML/DL 训练与推理前的标配步骤；本服务可运行。",
    },
    23: {
        "method": "PCA/MNF/ICA 将高度相关的上百波段压缩为低维主成分（常用 10–40 维）。",
        "value": "降冗余、抑噪声、降算力，加快训练与边缘部署。",
        "risk": "全波段直接训练易过拟合且耗时长，部署成本高。",
        "pipeline": "高维 Cube → 低维特征立方体 → 分类/回归",
        "related": "关联：#22 标准化、#24 波段选择、#34–37 分类",
        "tip": "本仓库部分小模型流水线内含 PCA；本服务可运行。",
    },
    24: {
        "method": "基于可分性准则、重要性评分或包装法，选出对任务区分力最强的波段子集。",
        "value": "在精度可接受前提下降低数据量，指导定波段传感器与轻量化部署。",
        "risk": "盲目全波段增加成本，未必提升精度。",
        "pipeline": "Cube +（可选）标签 → 优选波段列表 → 轻量模型/传感器配置",
        "related": "关联：#23 降维、#27–31 指数、特定作物分类",
        "tip": "波段选择强调「可解释子集」；PCA 强调「正交压缩」。",
    },
    25: {
        "method": "超像素/面向对象分割生成均质斑块，后续以对象为单位分类或统计。",
        "value": "减少椒盐噪声，结果更符合农田斑块连续性。",
        "risk": "纯像素分类易碎斑，地块边界可读性差。",
        "pipeline": "L2 影像 → 对象多边形 → 对象级分类或制图",
        "related": "关联：#34–38 分类、#44 后处理、#45 汇总",
        "tip": "对象尺度应与田块管理单元匹配。",
    },
    26: {
        "method": "按窗宽从立方体切取邻域 Patch 或光谱向量，并与标签对齐构成训练/推理样本。",
        "value": "为 CNN/Transformer 提供标准张量输入，统一训练与推理接口。",
        "risk": "缺少规范样本构建则深度学习管线无法稳定复现。",
        "pipeline": "Cube + 标签 + 窗宽 → 样本张量 (N,W,W,B)/(N,B) → #36–38",
        "related": "关联：#37 3D-CNN、#38 Transformer、本仓库 createImageCubes",
        "tip": "窗宽需在空间上下文与类别纯度之间权衡。",
    },
    27: {
        "method": "NDVI = (NIR − RED) / (NIR + RED)，输出约 −1～1 的单波段专题图。",
        "value": "最通用的植被绿度与光合活性表征，适合长势监测与物候观察。",
        "risk": "缺少快速长势产品时，监测业务缺乏标准化入口指标。",
        "pipeline": "L2 红光+近红外 → NDVI 图 → #45 地块均值/阈值告警",
        "related": "关联：#28 NDRE、#29 EVI/SAVI、#45 汇总",
        "tip": "茂密冠层易饱和，可并行使用 NDRE/EVI；本服务可运行。",
    },
    28: {
        "method": "NDRE = (NIR − RE) / (NIR + RE)，利用红边波段对叶绿素变化更敏感。",
        "value": "适用于密冠层与氮营养相关监测，发挥红边观测优势。",
        "risk": "仅用 NDVI 时，成熟期差异可能被饱和掩盖，施肥诊断不足。",
        "pipeline": "L2 红边+近红外 → NDRE 图 → 营养诊断专题 / #45",
        "related": "关联：#27 NDVI、#31 红边参数、#32 氮含量反演",
        "tip": "高光谱/多光谱红边配置的典型增值指标；本服务可运行。",
    },
    29: {
        "method": "EVI 引入蓝光抑制大气；SAVI/MSAVI 引入土壤调节因子，缓解土壤背景与饱和。",
        "value": "在苗期稀疏或过密冠层阶段提供比 NDVI 更稳妥的长势表征。",
        "risk": "生育期不匹配指数时，长势评估偏差增大。",
        "pipeline": "按生育期选择指数 → 专题图 → 地块汇总",
        "related": "关联：#27 NDVI、#28 NDRE",
        "tip": "茂密优先 EVI；稀疏苗期优先 SAVI/MSAVI。",
    },
    30: {
        "method": "NDMI≈(NIR−SWIR)/(NIR+SWIR)；NDWI/MNDWI 用绿/SWIR 等组合提取水分或水体。",
        "value": "支撑干旱、灌溉、湿地与农田积水等水分管理类专题。",
        "risk": "仅有绿度维度时，水分胁迫与淹田识别能力不足，农事建议不完整。",
        "pipeline": "L2 相关波段 → 水分/水体指数图 → 阈值告警 / #45",
        "related": "关联：#27 长势链、#42 异常、#43 变化",
        "tip": "与长势指数并列的第二条常用指标链；注意 SWIR 波段可用性。",
    },
    31: {
        "method": "从连续光谱提取红边位置（REP）、吸收谷深度、一阶导数特征等参数栅格。",
        "value": "体现高光谱相对宽波段多光谱的差异化监测能力（物候/胁迫早期）。",
        "risk": "仅做传统宽波段指数时，难以体现高光谱投入的增量价值。",
        "pipeline": "连续反射率 → 参数栅格 → 物候/胁迫专题",
        "related": "关联：#28 NDRE、#21 去包络、#32 反演",
        "tip": "适合作为高光谱能力说明与科研/精品监测产品。",
    },
    32: {
        "method": "以 PLS/随机森林/神经网络等建立光谱—生化量回归（叶绿素、氮、含水率等）。",
        "value": "将光谱映射为带物理单位的连续量，服务精准施肥与诊断。",
        "risk": "停留在定性看图，难以形成剂量级农事处方。",
        "pipeline": "光谱特征 + 地面化验样本 → 训练 → 连续量专题图",
        "related": "关联：#27–31 特征、#33 物理反演、#45 报表",
        "tip": "依赖代表性地面样本；无样本时不可盲目外推。",
    },
    33: {
        "method": "基于冠层辐射传输机理（如 PROSAIL 类）反演 LAI、叶绿素等参数。",
        "value": "可解释性强，适合高精度定量科研与高端业务。",
        "risk": "复杂场景仅靠经验回归可能外推失败；物理法成本更高。",
        "pipeline": "反射率 + 先验参数 → 物理反演 → 参数图（如 LAI）",
        "related": "关联：#13 大气校正、#32 经验回归",
        "tip": "与经验回归互补：机理约束 vs 数据驱动。",
    },
    34: {
        "method": "以光谱或降维特征训练 SVM/随机森林/逻辑回归，按像素输出类别；可计算 OA/AA/Kappa。",
        "value": "小样本条件下快速回答「地物/作物类别」问题，可解释性较好。",
        "risk": "缺少基线分类能力时，业务无法输出作物一张图，深度学习也缺少对照。",
        "pipeline": "特征 + 训练标签 → 分类图 → #44 后处理 → #45 面积",
        "related": "关联：#23 PCA、#37 深度学习、#44/#45",
        "tip": "本服务可运行，并返回精度指标便于验收。",
    },
    35: {
        "method": "光谱角制图（SAM）等：比较像元光谱与端元/光谱库的形状相似度。",
        "value": "在具备标准光谱库时，可直接进行目标/矿物等识别。",
        "risk": "库内物质无法快速识别，只能依赖大量标注分类。",
        "pipeline": "反射率 + 光谱库 → 匹配类别图/距离图",
        "related": "关联：#9 Smile 校正、#21 去包络、#41 解混",
        "tip": "地质与目标识别常用；一般作物分类更多用监督分类。",
    },
    36: {
        "method": "1D-CNN/RNN 沿光谱维自动提取特征并完成像素分类。",
        "value": "在光谱可分性强的场景中，优于手工指数与浅层模型。",
        "risk": "复杂空间结构场景下，仅光谱模型边界容易出错。",
        "pipeline": "单像素光谱 (B,) → 1D 网络 → 分类图 → #44/#45",
        "related": "关联：#22/#23 前处理、#37 空谱 CNN",
        "tip": "适合光谱区分度高、空间纹理较弱的任务。",
    },
    37: {
        "method": "2D/3D-CNN 同时利用邻域空间与光谱信息，输入邻域立方体 (W,W,B)。",
        "value": "农田精细分类与学术主流路线，通常可获得更高的边界与斑块精度。",
        "risk": "仅光谱分类在破碎地块与边界区域更容易产生误分。",
        "pipeline": "#26 Patch → 空谱网络 → 分类着色图 → #44/#45",
        "related": "关联：本仓库 hyper-spectral-small-modes 核心落点",
        "tip": "作物一张图、城市地物精细分类的优先技术方案。",
    },
    38: {
        "method": "Transformer 注意力或 GCN 图结构建模长程依赖与对象关系。",
        "value": "面向复杂地物与大场景的高精度升级路径。",
        "risk": "难分场景下传统 CNN 精度到顶时，需要更强关系建模能力。",
        "pipeline": "Patch/超像素图特征 → Transformer/GCN → 分类图",
        "related": "关联：#25 超像素、#37 CNN、#39 少样本",
        "tip": "算力与数据要求更高，适合研究型与高端产品。",
    },
    39: {
        "method": "迁移学习、元学习或少样本微调：源域预训练 + 目标域少量标注。",
        "value": "降低新区域外业认种与标注成本，加快业务扩展。",
        "risk": "每换场景都需大量标注时，规模化落地受限。",
        "pipeline": "源域模型 + 少量目标标签 → 目标场景分类图 → #45",
        "related": "关联：#34–38 分类家族、#26 样本",
        "tip": "衡量「可快速铺开」的关键能力指标。",
    },
    40: {
        "method": "语义分割/目标检测定位病斑、杂草等目标的边界或检测框，常融合 RGB。",
        "value": "支撑精准喷药与植保巡田，从「类别图」升级到「处置位置」。",
        "risk": "无法给出位置时只能粗放全田作业，农资与成本失控。",
        "pipeline": "Cube/RGB + 标注 → 掩膜/检测框/矢量 → 处方图",
        "related": "关联：#37 分割网络、#45 告警与报表",
        "tip": "输出 shp/geojson 便于对接植保设备。",
    },
    41: {
        "method": "端元提取 + 丰度反演（如 FCLS），估计混合像元内各类占地比例（0～1）。",
        "value": "在混种、稀疏植被或分辨率不足时提供亚像元占比信息。",
        "risk": "硬分类强迫唯一类别，导致比例估测失真。",
        "pipeline": "混合光谱 + 端元 → 丰度图 → 专题统计",
        "related": "关联：#35 光谱匹配、#34 硬分类（对照）",
        "tip": "输出连续丰度，而非单一类别标签。",
    },
    42: {
        "method": "基于统计或重建残差计算异常得分，找出光谱上偏离背景的像元。",
        "value": "无充分标注时用于病虫害爆发点、污染点等初筛告警。",
        "risk": "小范围异常被区域均值淹没，难以及时发现。",
        "pipeline": "单时相 Cube → 异常得分图 → 人工复核/补采样",
        "related": "关联：#5 云影（排除干扰）、#43 变化检测",
        "tip": "适合「先告警、再确认」的运维流程。",
    },
    43: {
        "method": "对配准后的多时相 L2/L3 做差值、比值或变化向量分析，提取变化区。",
        "value": "支撑灾损、砍伐、轮作、施工占地等「变化位置」监测。",
        "risk": "无变化产品时只能单期看图，难以回答阶段性变化问题。",
        "pipeline": "多时相配准 → 变化检测 → 变化掩膜/类型图 → 告警",
        "related": "关联：#11 辐射归一、#19 配准、#45 告警",
        "tip": "配准与辐射一致性是前提，否则伪变化显著。",
    },
    44: {
        "method": "形态学滤波、CRF 或小斑剔除，去除椒盐噪声与过小斑块。",
        "value": "提升分类图连续性与可读性，稳定后续面积统计。",
        "risk": "碎斑过多影响验收观感，并干扰地块面积汇总。",
        "pipeline": "原始 LabelMap → 平滑分类图 → #45 地块汇总",
        "related": "关联：#34–38 分类、#25 对象、#45 汇总",
        "tip": "出图验收前的标准后处理步骤。",
    },
    45: {
        "method": "以地块矢量为单元，对分类/指数栅格做分区统计，结合阈值生成报表与告警。",
        "value": "将像素结果转化为面积、占比、告警列表等可对接业务系统的决策产品。",
        "risk": "缺少汇总时结果停留在技术图，难以进入业务系统与管理流程。",
        "pipeline": "L3 产品 + 地块边界 + 规则 → 面积表/JSON/专题图/告警/API",
        "related": "关联：全链路出口；本服务可运行",
        "tip": "L4 核心交付形态：报表 + 专题图 + 接口。",
    },
}


def rgb(t):
    return RGBColor(*t)


def set_run_font(run, size=14, bold=False, color=None, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color or C["ink"])
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def write_para(p, text, size=14, bold=False, color=None, align=None, space_after=4):
    p.clear()
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color)
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    if align is not None:
        p.alignment = align


def add_text(slide, left, top, width, height, lines, *, size=14, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, sz, bd, col = item
            write_para(p, text, size=sz, bold=bd, color=col, align=align, space_after=3)
        else:
            write_para(p, item, size=size, bold=bold, color=color, align=align, space_after=3)
    return box


def add_rect(slide, left, top, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = rgb(line)
    return sh


def add_round(slide, left, top, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = rgb(line)
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    return sh


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, 13.333, 7.5, C["bg"])
    return s


def footer(slide, page: int, total: int, section: str = ""):
    add_rect(slide, 0, 7.18, 13.333, 0.32, C["primary_dark"])
    add_text(
        slide, 0.35, 7.2, 9, 0.28,
        f"高光谱算法体系  ·  {section}" if section else "高光谱算法体系",
        size=11, color=C["white"],
    )
    add_text(slide, 10.5, 7.2, 2.5, 0.28, f"{page} / {total}", size=11, color=C["white"], align=PP_ALIGN.RIGHT)


def title_header(slide, title: str, subtitle: str = "", pill: str = "", pill_fill=None):
    add_rect(slide, 0, 0, 13.333, 0.92, C["primary_dark"])
    add_rect(slide, 0, 0.92, 13.333, 0.05, C["accent"])
    add_text(slide, 0.35, 0.14, 10.6, 0.45, title, size=22, bold=True, color=C["white"])
    if subtitle:
        add_text(slide, 0.35, 0.55, 10.6, 0.3, subtitle, size=11, color=(167, 243, 208))
    if pill:
        add_round(slide, 10.85, 0.22, 2.1, 0.42, pill_fill or C["warm"])
        add_text(slide, 10.85, 0.26, 2.1, 0.35, pill, size=11, bold=True, color=C["white"], align=PP_ALIGN.CENTER)


def section_card(slide, left, top, w, h, title, body_lines, accent=None, title_size=12, body_size=11):
    accent = accent or C["primary"]
    add_round(slide, left, top, w, h, C["card"], C["line"])
    add_rect(slide, left, top, 0.07, h, accent)
    lines = [(title, title_size, True, accent)]
    for b in body_lines:
        if b == "":
            lines.append(("", 4, False, C["ink"]))
        else:
            lines.append((b, body_size, False, C["ink"]))
    add_text(slide, left + 0.16, top + 0.08, w - 0.28, h - 0.14, lines)


def parse_algorithms(md_text: str) -> list[dict]:
    items = []
    headings = list(re.finditer(r"^####\s+(\d+)\.\s+(.+?)\s*$", md_text, re.M))
    for index, match in enumerate(headings):
        block_end = headings[index + 1].start() if index + 1 < len(headings) else len(md_text)
        block = md_text[match.end():block_end]
        values = {}
        for label, key in (("作用", "role"), ("使用场景", "scene"), ("数据输入", "inp"), ("数据输出", "out")):
            field = re.search(rf"^\|\s*\*\*{label}\*\*\s*\|\s*(.*?)\s*\|\s*$", block, re.M)
            values[key] = field.group(1).strip().replace("**", "") if field else ""
        items.append({"num": int(match.group(1)), "name": match.group(2).strip(), **values})
    level_map = {}
    for line in md_text.splitlines():
        m = re.match(r"\|\s*(\d+)\s*\|\s*([^|]+)\s*\|\s*([^|]+)\s*\|\s*([^|]+)\s*\|", line)
        if m and m.group(1).isdigit():
            level_map[int(m.group(1))] = m.group(2).strip()
    # 正式表述：避免口语化「领导」用词
    formal_replace = {
        "领导汇报": "业务汇报",
        "领导、客户": "业务方与客户",
        "领导/客户": "业务方与客户",
        "领导": "业务决策",
    }
    for it in items:
        it["level"] = level_map.get(it["num"], "")
        for key in ("role", "scene", "inp", "out"):
            text = it[key]
            for a, b in formal_replace.items():
                text = text.replace(a, b)
            it[key] = text
        it.update(ENRICH.get(it["num"], {}))
    return items


def validate_content(algos: list[dict]) -> None:
    """校验算法编号、接口标识和最新实现状态口径。"""
    numbers = [item["num"] for item in algos]
    assert numbers == list(range(1, 46)), f"算法编号异常：{numbers}"
    assert len(ALG_IDS) == len(set(ALG_IDS)) == 45, "algorithm_id 缺失或重复"
    assert len(IMPL) == 12, f"可运行算法数量异常：{len(IMPL)}"
    assert len(set(range(1, 46)) - IMPL) == 33, "骨架算法数量异常"
    print("内容校验通过：45 项；可运行 12；骨架 33")


def build():
    md = LIST_MD.read_text(encoding="utf-8")
    algos = parse_algorithms(md)
    validate_content(algos)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slides_meta = []

    def track(s, section):
        slides_meta.append((s, section))
        return s

    # 封面
    s = blank(prs)
    track(s, "")
    add_rect(s, 0, 0, 13.333, 7.5, C["primary_dark"])
    add_rect(s, 0, 0, 0.28, 7.5, C["accent"])
    add_text(s, 0.8, 1.7, 11.5, 0.7, "高光谱数据处理与算法体系", size=34, bold=True, color=C["white"])
    add_text(
        s, 0.8, 2.55, 11.5, 1.2,
        [
            ("架构设计  ·  45 项算法详解  ·  服务能力与对接", 18, False, (167, 243, 208)),
            ("政企科技蓝统一版｜专业培训与业务汇报材料", 14, False, (171, 207, 232)),
        ],
    )
    add_round(s, 0.8, 4.0, 4.0, 0.5, C["accent"])
    add_text(s, 0.8, 4.08, 4.0, 0.4, "L0 → L1 → L2 → L3 → L4", size=13, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    add_text(
        s, 0.8, 4.9, 11.5, 1.4,
        [
            ("核心链路：原始 DN → 反射率正射立方体 → 分类/指数产品 → 面积报表与业务告警", 13, False, (226, 232, 240)),
            ("内容口径：45 项能力覆盖 L0–L4；12 项可运行，33 项为接口骨架", 12, False, (171, 207, 232)),
            ("配套服务：统一 HTTP API；GeoTIFF / GeoJSON / CSV 标准格式", 12, False, (171, 207, 232)),
        ],
    )

    # 目录
    s = blank(prs)
    track(s, "目录")
    title_header(s, "目录", "从采集到业务决策的完整能力说明")
    cards = [
        ("01  架构设计", ["数据层级 L0–L4 定义与用户", "分类链与长势指标链", "业界交换文件格式", "端到端业务处理路径"], C["primary"]),
        ("02  算法清单（45项）", ["按层级说明功能与方法要点", "明确输入输出与产品形态", "业务价值、缺失影响与流程位置", "关联算法与实施提示"], C["accent"]),
        ("03  服务能力与对接", ["统一 HTTP 服务定位", "已实现能力与接口契约", "联调步骤与调用示例", "能力边界与后续建设"], C["warm"]),
    ]
    for i, (t, bullets, col) in enumerate(cards):
        x = 0.4 + i * 4.25
        add_round(s, x, 1.25, 4.05, 5.5, C["card"], C["line"])
        add_rect(s, x, 1.25, 4.05, 0.65, col)
        add_text(s, x + 0.2, 1.35, 3.7, 0.45, t, size=15, bold=True, color=C["white"])
        add_text(s, x + 0.25, 2.2, 3.55, 4.2, [(f"• {b}", 13, False, C["ink"]) for b in bullets])

    # 第一部分
    s = blank(prs)
    track(s, "架构")
    title_header(s, "第一部分｜架构设计", "先建立数据产品分层，再理解算法落点")
    section_card(s, 0.35, 1.2, 6.2, 5.6, "核心结论", [
        "高光谱业务不是「一组光谱直接得到一个指标」。",
        "业界将链路拆分为多层数据产品，层层增值：",
        "",
        "L0  原始 DN 计数 + POS/元数据",
        "L1  辐亮度（具备物理单位）",
        "L2  反射率正射立方体（可上地图）",
        "L3  分类图 / 指数图 / 反演图",
        "L4  面积、报表、告警、业务 API",
        "",
        "本仓库分类小模型主要落在 L2→L3；",
        "业务系统最终消费的多为 L4 产品。",
    ], C["primary"], 14, 12)
    section_card(s, 6.8, 1.2, 6.15, 5.6, "分层价值（摘要）", [
        "L0：保住可追溯的原始观测",
        "L1：进入可定量的物理量空间",
        "L2：形成算法可消费的标准底图",
        "L3：提取类别、长势、胁迫等信息",
        "L4：转化为可管理、可对接的结论",
        "",
        "两条并行分析链均建立在 L2 之上：",
        "① 分类链：回答「种了什么、各多少」",
        "② 指标链：回答「长势/水分如何」",
        "",
        "培训重点：先记住分层，再记算法名称。",
    ], C["accent"], 14, 12)

    s = blank(prs)
    track(s, "架构")
    title_header(s, "数据层级 L0–L4", "形态 · 主要用户 · 可否直接用于管理决策")
    levels = [
        ("L0 原始层", ["形态：DN 帧/条带", "POS、日志、定标板", "", "用户：采集与预处理", "", "决策：否", "价值：原始可追溯"]),
        ("L1 传感器产品", ["形态：辐亮度", "可含粗几何", "", "用户：处理中心", "", "决策：否", "价值：物理可比"]),
        ("L2 分析底图", ["形态：反射率", "正射立方体", "", "用户：算法研发", "", "决策：专业预览", "价值：标准输入"]),
        ("L3 信息产品", ["形态：分类/指数", "反演/矢量", "", "用户：分析与产品", "", "决策：部分可读", "价值：信息提取"]),
        ("L4 业务决策", ["形态：亩数报表", "告警与 API", "", "用户：业务与系统", "", "决策：是", "价值：管理闭环"]),
    ]
    for i, (title, body) in enumerate(levels):
        x = 0.3 + i * 2.58
        add_round(s, x, 1.2, 2.48, 5.6, C["card"], C["line"])
        add_rect(s, x, 1.2, 2.48, 0.55, C["primary"] if i < 4 else C["warm"])
        add_text(s, x + 0.08, 1.28, 2.3, 0.4, title, size=12, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_text(s, x + 0.12, 1.95, 2.25, 4.6, [(b, 11, False, C["ink"]) for b in body])

    s = blank(prs)
    track(s, "架构")
    title_header(s, "端到端处理路径（以作物分类为例）", "从飞行采集到业务系统可读结果")
    steps = [
        ("① 采集", "无人机高光谱（可选 RGB）飞行；布置白板；记录样点类别与地块边界。"),
        ("② 预处理", "暗电流/坏线/定标 → 正射 → 镶嵌，形成 L2 反射率正射立方体（GeoTIFF）。"),
        ("③ 分析", "样本标注 + 分类模型（传统 ML / 空谱 CNN 等）→ L3 分类着色图。"),
        ("④ 后处理", "平滑与小斑剔除，提升图斑连续性，稳定面积统计。"),
        ("⑤ 汇总", "按地块分区统计亩数与占比 → L4 报表 / JSON / 专题图。"),
        ("⑥ 应用", "业务系统、GIS、大屏消费 L4；并行可走 NDVI/NDRE 长势监测链。"),
    ]
    for i, (h, t) in enumerate(steps):
        y = 1.15 + i * 0.92
        add_round(s, 0.35, y, 12.6, 0.82, C["card"], C["line"])
        add_rect(s, 0.35, y, 1.55, 0.82, C["primary"] if i < 5 else C["warm"])
        add_text(s, 0.4, y + 0.22, 1.45, 0.4, h, size=13, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_text(s, 2.15, y + 0.22, 10.5, 0.45, t, size=13, color=C["ink"])

    s = blank(prs)
    track(s, "架构")
    title_header(s, "两条业务主链路", "均以 L2 反射率正射立方体为共同输入")
    section_card(s, 0.35, 1.2, 6.2, 5.6, "A. 作物/地物分类链", [
        "L2 反射率正射立方体",
        "→ 样本构建与特征处理（#20–26）",
        "→ 分类模型（#34–39）",
        "→ 分类后处理（#44）",
        "→ 地块面积与专题图（#45）",
        "",
        "回答的问题：",
        "· 各地块种植了什么？",
        "· 各类面积与空间分布如何？",
        "",
        "本仓库小模型主要服务本链路。",
        "验收关注：OA/Kappa + 面积核对。",
    ], C["primary"], 13, 12)
    section_card(s, 6.8, 1.2, 6.15, 5.6, "B. 长势/胁迫指标链", [
        "L2 反射率正射立方体",
        "→ NDVI/NDRE 等指数（#27–31）",
        "→ 可选生化量反演（#32–33）",
        "→ 热力专题图",
        "→ 地块均值与阈值告警（#45）",
        "",
        "回答的问题：",
        "· 长势空间差异如何？",
        "· 是否存在干旱/缺肥风险区？",
        "",
        "链路更轻、更适合高频监测。",
        "验收关注：阈值规则与告警闭环。",
    ], C["warm"], 13, 12)

    s = blank(prs)
    track(s, "架构")
    title_header(s, "业界文件格式介绍", "服务与业务系统按标准格式交换，而非教学用 npy")
    rows = [
        ("GeoTIFF (.tif)", "多波段栅格 + 地理参考；主格式", "L2 立方体、指数图、分类图"),
        ("ENVI", "经典高光谱交换；服务侧可读", "处理中心存量数据接入"),
        ("GeoJSON", "测区、地块、斑块矢量", "分区统计、AOI 裁剪"),
        ("CSV", "POS 轨迹、光谱库、端元表", "几何与光谱匹配辅助"),
        ("JSON 报表", "统计指标 + 产物路径 + 告警", "业务系统 / 大屏直接消费"),
    ]
    for i, (a, b, c) in enumerate(rows):
        y = 1.15 + i * 1.1
        add_round(s, 0.35, y, 12.6, 0.98, C["card"], C["line"])
        add_rect(s, 0.35, y, 0.1, 0.98, C["accent"])
        add_text(s, 0.65, y + 0.28, 2.8, 0.4, a, size=13, bold=True, color=C["primary"])
        add_text(s, 3.6, y + 0.18, 5.3, 0.65, b, size=12, color=C["ink"])
        add_text(s, 9.1, y + 0.18, 3.6, 0.65, f"用途：{c}", size=12, color=C["muted"])

    # 第二部分引导
    s = blank(prs)
    track(s, "算法清单")
    title_header(s, "第二部分｜算法清单（45 项）", "每项说明：功能 · 方法 · 场景 · 输入输出 · 价值 · 影响 · 流程")
    section_card(s, 0.35, 1.2, 12.6, 5.6, "阅读与培训建议", [
        "1. 先按层级理解「本层解决什么业务问题」，再记忆算法名称。",
        "2. 每页包含：功能定义、方法要点、适用场景、输入输出、业务价值、缺失影响、流程位置与关联。",
        "3. 标注「可运行」表示当前 HTTP 服务可联调（教学实现）；其余为业界能力地图，接口已预留。",
        "4. 分层覆盖：",
        "   #1–5 采集与质检　#6–11 传感器产品　#12–26 分析底图与特征",
        "   #27–43 信息提取　#44–45 出图与业务汇总",
        "5. 与仓库关系：L0–L2 工程能力多在处理中心；本仓库强项在 L2→L3 分类及可运行教学 API。",
    ], C["primary"], 14, 13)

    groups = [
        ("L0 及上机前", "1–5", "飞行与原始数据质量控制", list(range(1, 6))),
        ("L0→L1 传感器产品", "6–11", "从 DN 到物理辐亮度", list(range(6, 12))),
        ("L1→L2 分析底图（上）", "12–19", "反射率、几何、镶嵌、配准", list(range(12, 20))),
        ("L2 分析前处理", "20–26", "清洗、特征、样本构建", list(range(20, 27))),
        ("L3 指数与反演", "27–33", "长势、水分、生化量", list(range(27, 34))),
        ("L3 分类与检测", "34–43", "类别、斑块、异常、变化", list(range(34, 44))),
        ("L3→L4 业务出口", "44–45", "后处理与决策产品", list(range(44, 46))),
    ]
    for page_groups in (groups[:4], groups[4:]):
        s = blank(prs)
        track(s, "算法清单")
        title_header(s, "算法分层总览", "45 项能力与业务问题的对应关系")
        for i, (title, rng, desc, ids) in enumerate(page_groups):
            y = 1.15 + i * 1.4
            names = "、".join(
                re.sub(r"（.*?）|\(.*?\)", "", next(a["name"] for a in algos if a["num"] == n))[:12]
                for n in ids
            )
            add_round(s, 0.35, y, 12.6, 1.25, C["card"], C["line"])
            add_rect(s, 0.35, y, 2.5, 1.25, C["primary"])
            add_text(s, 0.45, y + 0.25, 2.3, 0.4, title, size=12, bold=True, color=C["white"])
            add_text(s, 0.45, y + 0.7, 2.3, 0.35, f"#{rng}", size=11, color=(167, 243, 208))
            add_text(s, 3.1, y + 0.2, 9.5, 0.35, desc, size=14, bold=True, color=C["primary"])
            add_text(s, 3.1, y + 0.6, 9.5, 0.5, names, size=11, color=C["muted"])

    # 每个算法一页：统一五区结构，兼顾领导汇报与技术培训
    for a in algos:
        s = blank(prs)
        track(s, "算法清单")
        num = a["num"]
        aid = ALG_IDS[num - 1]
        is_impl = num in IMPL
        status = "可运行｜教学级实现" if is_impl else "骨架｜核心算法待实现"
        title_header(
            s,
            f"{num:02d}. {a['name']}",
            f"层级：{a.get('level', '')}　｜　algorithm_id：{aid}",
            pill=status,
            pill_fill=C["ok"] if is_impl else C["warm"],
        )

        # 左侧主叙事：问题、原理、输入输出
        section_card(s, 0.3, 1.12, 8.15, 1.48, "01｜解决问题", [
            a["role"],
            f"典型场景：{a['scene']}",
        ], C["primary"], 13, 12)

        section_card(s, 0.3, 2.75, 8.15, 1.72, "02｜技术原理", [
            a.get("method", "—"),
            f"工程提示：{a.get('tip', '按项目验收标准设置阈值与抽检。')}",
        ], C["accent"], 13, 12)

        section_card(s, 0.3, 4.62, 8.15, 2.18, "03｜输入与输出", [
            f"输入：{a['inp']}",
            "",
            f"输出：{a['out']}",
            "",
            f"流程位置：{a.get('pipeline', '—')}",
        ], C["primary"], 13, 11)

        # 右侧决策信息：价值与状态
        section_card(s, 8.65, 1.12, 4.38, 2.55, "04｜业务价值", [
            a.get("value", "—"),
            "",
            f"风险提示：{a.get('risk', '—')}",
        ], C["ok"], 13, 11)

        state_lines = [
            f"当前状态：{status}",
            "",
            "能力说明：" + (
                "当前接口具备可运行示例，生产上线仍需真实数据精度、性能与异常恢复验收。"
                if is_impl else
                "当前仅提供统一目录与接口契约，不代表核心算法能力已完成。"
            ),
            "",
            a.get("related", "—"),
            "",
            f"API：/api/v1/{aid}/run",
        ]
        section_card(s, 8.65, 3.82, 4.38, 2.98, "05｜状态与工程边界", state_lines, C["accent"] if is_impl else C["warm"], 13, 11)

    # 第三部分
    s = blank(prs)
    track(s, "服务对接")
    title_header(s, "第三部分｜服务能力与对接", "统一 HTTP 入口，业界文件格式入出")
    section_card(s, 0.35, 1.2, 12.6, 5.6, "服务定位与解决的问题", [
        "定位：高光谱算法能力的统一 HTTP 服务（单进程、一算法一目录）。",
        "",
        "解决的问题：",
        "· 业务系统无需各自实现遥感算法，通过标准文件与接口调用即可。",
        "· 培训与联调统一入口：清单编号 ↔ 目录 ↔ API 一一对应。",
        "· 先以教学实现打通闭环，再按优先级将骨架算法补齐到生产级。",
        "",
        "调用形态：上传 GeoTIFF / GeoJSON / CSV → 返回 JSON（含统计与产物路径）。",
        "交互文档：http://服务地址:28800/docs",
    ], C["primary"], 14, 13)

    s = blank(prs)
    track(s, "服务对接")
    title_header(s, "对不同角色的价值", "同一套能力，面向业务、产品与算法三类读者")
    section_card(s, 0.3, 1.15, 4.15, 5.65, "业务与管理", [
        "将光谱观测转化为分类图、长势图、面积与告警。",
        "支撑农情监测、作物一张图等应用场景。",
        "降低对单一算法专家的路径依赖。",
        "",
        "验收关注：",
        "· 专题图可读性",
        "· 面积可核对",
        "· 告警可执行",
    ], C["warm"], 13, 12)
    section_card(s, 4.6, 1.15, 4.15, 5.65, "产品与应用研发", [
        "统一 API，可按场景组合算法。",
        "业界格式便于接入 GIS / App。",
        "每算法独立目录、README 与 testdata。",
        "可从骨架渐进演进到生产实现。",
        "",
        "对接关注：",
        "· 契约稳定",
        "· OpenAPI 可调试",
        "· 样例可一键跑通",
    ], C["primary"], 13, 12)
    section_card(s, 8.9, 1.15, 4.1, 5.65, "算法与数据工程", [
        "对齐 L0–L4 业界分层语言。",
        "与小模型仓库能力可衔接。",
        "输入输出契约清晰，便于验收。",
        "",
        "实现关注：",
        "· implemented 标记",
        "· 产物 files 路径",
        "· 与清单编号一致",
    ], C["accent"], 13, 12)

    s = blank(prs)
    track(s, "服务对接")
    title_header(s, "当前可运行能力（可联调）", "12 项教学级实现；33 项统一接口骨架")
    impl_lines = [f"#{n:02d} {next(a['name'] for a in algos if a['num']==n)}" for n in sorted(IMPL)]
    section_card(s, 0.35, 1.2, 7.65, 5.6, "可运行 12 项", [
        *impl_lines,
    ], C["ok"], 13, 11)
    section_card(s, 8.2, 1.2, 4.75, 2.45, "能力覆盖", [
        "反射率定标与立方体清洗",
        "特征处理与植被指数",
        "传统分类与空谱 CNN",
        "检测、异常与地块汇总",
    ], C["primary"], 13, 12)
    section_card(s, 8.2, 3.85, 4.75, 2.95, "工程边界", [
        "33 项为统一接口骨架：",
        "目录与 POST /api/v1/{id}/run 已预留，",
        "implemented=false，核心算法需按场景补齐。",
        "",
        "可运行项仍需通过真实数据精度、性能、",
        "稳定性和安全验收后方可生产上线。",
        "",
        "详见：docs/算法API测试清单.md",
    ], C["warm"], 13, 11)

    s = blank(prs)
    track(s, "服务对接")
    title_header(s, "对接步骤与调用示例", "工作目录：algorithm/source")
    section_card(s, 0.3, 1.15, 6.3, 5.65, "三步对接", [
        "1）启动服务",
        "./scripts/start.sh",
        "默认：http://127.0.0.1:28800",
        "",
        "2）查询清单",
        "GET /api/v1/algorithms",
        "",
        "3）调用算法",
        "POST /api/v1/{algorithm_id}/run",
        "-F file=@xxx.tif",
        "-F file2=@yyy.geojson（可选）",
        "-F params='{...}'",
        "",
        "返回：success / data / files",
    ], C["primary"], 13, 12)
    section_card(s, 6.8, 1.15, 6.15, 5.65, "示例：NDVI（#27）", [
        "curl -X POST \\",
        "  \"http://127.0.0.1:28800/api/v1/27_ndvi/run\" \\",
        "  -F \"file=@algorithms/27_ndvi/testdata/input.tif\" \\",
        "  -F 'params={\"red_band\":2,\"nir_band\":3}'",
        "",
        "业务侧可消费：",
        "· data.mean / min / max",
        "· files.ndvi_tif（GIS 底图）",
        "· files.preview_png（预览）",
        "",
        "同类方式对接 #34 分类、#45 汇总。",
    ], C["accent"], 13, 11)

    s = blank(prs)
    track(s, "服务对接")
    title_header(s, "能力边界与建设路径", "汇报与立项需明确范围，避免预期错位")
    points = [
        ("当前范围", "12 项教学/联调级实现 + 33 项接口骨架 + 完整业界流程能力地图。"),
        ("非当前范围", "完整生产级正射、大气、SfM 等通常由专用处理软件/处理中心承担。"),
        ("仓库强项", "L2→L3 分类小模型（hyper-spectral-small-modes）可继续接入本服务。"),
        ("建议路径", "先用 12 项可运行接口打通业务闭环，再按场景优先级将 33 项骨架补到生产级。"),
    ]
    for i, (h, t) in enumerate(points):
        y = 1.2 + i * 1.3
        add_round(s, 0.35, y, 12.6, 1.15, C["card"], C["line"])
        add_rect(s, 0.35, y, 2.3, 1.15, C["primary"] if i != 1 else C["warm"])
        add_text(s, 0.45, y + 0.35, 2.1, 0.45, h, size=13, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_text(s, 2.9, y + 0.35, 9.7, 0.5, t, size=13, color=C["ink"])

    s = blank(prs)
    track(s, "总结")
    title_header(s, "总结与下一步", "")
    section_card(s, 0.35, 1.2, 4.0, 2.35, "01｜架构体系", [
        "L0→L4 数据产品分层清晰。",
        "分类链与指标链共享 L2 标准底图，",
        "最终在 L4 形成业务决策闭环。",
    ], C["primary"], 13, 12)
    section_card(s, 4.65, 1.2, 4.0, 2.35, "02｜算法能力", [
        "45 项能力覆盖采集、预处理、",
        "特征、分类、检测与业务汇总。",
        "12 项可运行，33 项为接口骨架。",
    ], C["accent"], 13, 12)
    section_card(s, 8.95, 1.2, 4.0, 2.35, "03｜服务对接", [
        "统一 HTTP API 与业界文件格式。",
        "GeoTIFF / GeoJSON / CSV 入出，",
        "支持培训、联调与渐进式建设。",
    ], C["ok"], 13, 12)
    section_card(s, 0.35, 3.85, 12.6, 2.95, "下一步建设建议", [
        "① 选定 1–2 个业务场景（作物分类 / 长势监测），形成端到端应用演示。",
        "② 按场景价值与交付风险排序，将关键骨架算法补齐为可验收的生产实现。",
        "③ 接入分类小模型与模型版本管理，建立数据—算法—业务反馈闭环。",
        "④ 补齐鉴权、异步任务、日志监控、精度基线与产物生命周期治理。",
    ], C["warm"], 13, 12)

    s = blank(prs)
    track(s, "")
    add_rect(s, 0, 0, 13.333, 7.5, C["primary_dark"])
    add_rect(s, 0, 0, 0.28, 7.5, C["accent"])
    add_text(s, 0.8, 2.7, 11, 0.7, "谢谢", size=40, bold=True, color=C["white"])
    add_text(
        s, 0.8, 3.6, 11, 1.5,
        [
            ("资料：algorithm/docs/（清单、流程、测试清单、本 PPT）", 15, False, (203, 213, 225)),
            ("服务：algorithm/source/（HTTP API · 端口 28800）", 15, False, (203, 213, 225)),
        ],
    )

    total = len(slides_meta)
    for i, (slide, section) in enumerate(slides_meta, 1):
        if i in (1, total):
            continue
        footer(slide, i, total, section)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print("wrote", OUT, "slides", len(prs.slides), "algos", len(algos))


def main():
    parser = argparse.ArgumentParser(description="生成高光谱算法体系 v4 PPT")
    parser.add_argument("--validate-only", action="store_true", help="仅校验内容与状态口径")
    args = parser.parse_args()
    if args.validate_only:
        md = LIST_MD.read_text(encoding="utf-8")
        validate_content(parse_algorithms(md))
        return
    build()


if __name__ == "__main__":
    main()
